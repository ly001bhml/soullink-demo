
SOULLINK_THERAPY_SUFFIX = '【心理陪伴回复要求】\n你现在要按心理陪伴型数字人的方式回复。先接住用户情绪，用一句自然的话说出你听到的感受；再结合用户说的具体内容，给出一个很小、能立刻做的建议；最后用一个温和的问题引导用户继续说。回复不要像心理学教材，不要诊断用户，不要说教，不要使用“你必须、你应该、别想太多、坚强点”。保持卡通角色的口吻，但心理场景优先稳定、真诚、具体。一般回复2到4句，适合直接播报。'

# -*- coding: utf-8 -*-
import os
import json
import time
import threading
import requests
import datetime
import schedule
import textwrap
from dataclasses import dataclass
from typing import Any, Callable, Dict, List, Literal, Optional, TypedDict, Tuple
from collections.abc import Mapping, Sequence
from langchain_openai import ChatOpenAI
from langchain_core.messages import HumanMessage, SystemMessage
from langgraph.graph import END, START, StateGraph

# 新增：本地知识库相关导入
import re
from pathlib import Path
import docx
from docx.document import Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# 用于处理 .doc 文件的库
try:
    import win32com.client
    WIN32COM_AVAILABLE = True
except ImportError:
    WIN32COM_AVAILABLE = False

from utils import util
import utils.config_util as cfg
from genagents.genagents import GenerativeAgent
from genagents.modules.memory_stream import ConceptNode
from urllib3.exceptions import InsecureRequestWarning
from scheduler.thread_manager import MyThread
from core import content_db
from core import stream_manager
from faymcp import tool_registry as mcp_tool_registry
# os.environ["LANGCHAIN_TRACING_V2"] = "true"
# os.environ["LANGCHAIN_ENDPOINT"] = "https://api.smith.langchain.com"
# os.environ["LANGCHAIN_API_KEY"] = "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"
# os.environ["LANGCHAIN_PROJECT"] = "fay3.9.1_github"

# 加载配置
cfg.load_config()

# 禁用不安全请求警告
requests.packages.urllib3.disable_warnings(category=InsecureRequestWarning)

agents = {}  # type: dict[str, GenerativeAgent]
agent_lock = threading.RLock()  # 使用可重入锁保护agent对象
reflection_lock = threading.RLock()  # 使用可重入锁保护reflection_time
save_lock = threading.RLock()  # 使用可重入锁保护save_time
reflection_time = None
save_time = None
motivational_quote_lock = threading.RLock()
motivational_quote_cooldowns: Dict[str, float] = {}
MOTIVATIONAL_QUOTE_MIN_INTERVAL_SECONDS = 15 * 60

memory_cleared = False  # 添加记忆清除标记
# 新增: 当前会话用户名及按用户获取memory目录的辅助函数
current_username = None  # 当前会话用户名

llm = ChatOpenAI(
        model=cfg.gpt_model_engine,
        base_url=cfg.gpt_base_url,
        api_key=cfg.key_gpt_api_key,
        streaming=True,
        max_tokens=10000 , # 限制回复长度，约50字内，朗读20秒内
        extra_body={"enable_thinking": False}
    )


@dataclass
class WorkflowToolSpec:
    name: str
    description: str
    schema: Dict[str, Any]
    executor: Callable[[Dict[str, Any], int], Tuple[bool, Optional[str], Optional[str]]]
    example_args: Dict[str, Any]


class ToolCall(TypedDict):
    name: str
    args: Dict[str, Any]


class ToolResult(TypedDict, total=False):
    call: ToolCall
    success: bool
    output: Optional[str]
    error: Optional[str]
    attempt: int


class ConversationMessage(TypedDict):
    role: Literal["user", "assistant"]
    content: str


class AgentState(TypedDict, total=False):
    request: str
    messages: List[ConversationMessage]
    tool_results: List[ToolResult]
    next_action: Optional[ToolCall]
    status: Literal["planning", "needs_tool", "completed", "failed"]
    final_response: Optional[str]
    final_messages: Optional[List[SystemMessage | HumanMessage]]
    planner_preview: Optional[str]
    audit_log: List[str]
    context: Dict[str, Any]
    error: Optional[str]
    max_steps: int


weather_schema = {
    "type": "object",
    "properties": {
        "location": {"type": "string", "description": "要查询天气的城市或地区，例如 济南、上海、北京"},
    },
    "required": ["location"],
}


def _truncate_text(text: Any, limit: int = 400) -> str:
    text_str = "" if text is None else str(text)
    if len(text_str) <= limit:
        return text_str
    return text_str[:limit] + "..."


_SMALL_TALK_PATTERNS = [
    re.compile(pattern)
    for pattern in [
        r"^(你|您)?好[呀啊呢哦哈~！!。.\s]*$",
        r"^嗨[呀啊呢哦哈~！!。.\s]*$",
        r"^hi[~!.\s]*$",
        r"^hello[~!.\s]*$",
        r"^早上好[呀啊呢哦~！!。.\s]*$",
        r"^中午好[呀啊呢哦~！!。.\s]*$",
        r"^下午好[呀啊呢哦~！!。.\s]*$",
        r"^晚上好[呀啊呢哦~！!。.\s]*$",
    ]
]


def _looks_like_small_talk(request: str) -> bool:
    text = (request or "").strip().lower()
    if not text or len(text) > 12:
        return False
    return any(pattern.fullmatch(text) for pattern in _SMALL_TALK_PATTERNS)


def _should_skip_workflow_for_request(
    request: str,
    conversation: List["ConversationMessage"],
) -> bool:
    if not _looks_like_small_talk(request):
        return False

    normalized_request = (request or "").strip()
    prior_messages = [
        message for message in (conversation or [])
        if (message.get("content") or "").strip() != normalized_request
    ]
    return len(prior_messages) == 0


def _extract_text_from_result(value: Any, *, depth: int = 0) -> List[str]:
    """Try to pull human-readable text snippets from tool results."""
    if value is None:
        return []
    if depth > 6:
        return []
    if isinstance(value, (str, int, float, bool)):
        text = str(value).strip()
        return [text] if text else []
    if isinstance(value, Mapping):
        # Prefer explicit text/content fields
        if "text" in value and not isinstance(value["text"], (dict, list, tuple)):
            text = str(value["text"]).strip()
            return [text] if text else []
        if "content" in value:
            segments: List[str] = []
            for item in value.get("content", []):
                segments.extend(_extract_text_from_result(item, depth=depth + 1))
            if segments:
                return segments
        segments = []
        for key, item in value.items():
            if key in {"meta", "annotations", "uid", "id", "messageId"}:
                continue
            segments.extend(_extract_text_from_result(item, depth=depth + 1))
        return segments
    if isinstance(value, Sequence) and not isinstance(value, (bytes, bytearray)):
        segments: List[str] = []
        for item in value:
            segments.extend(_extract_text_from_result(item, depth=depth + 1))
        return segments
    if hasattr(value, "text") and not callable(getattr(value, "text")):
        text = str(getattr(value, "text", "")).strip()
        return [text] if text else []
    if hasattr(value, "__dict__"):
        return _extract_text_from_result(vars(value), depth=depth + 1)
    text = str(value).strip()
    return [text] if text else []


def _normalize_tool_output(result: Any) -> str:
    """Convert structured tool output to a concise human-readable string."""
    if result is None:
        return ""
    segments = _extract_text_from_result(result)
    if segments:
        cleaned = [segment for segment in segments if segment]
        if cleaned:
            return "\n".join(dict.fromkeys(cleaned))
    try:
        return json.dumps(result, ensure_ascii=False, default=lambda o: getattr(o, "__dict__", str(o)))
    except TypeError:
        return str(result)


def _truncate_history(history: List[ToolResult], limit: int = 6) -> str:
    if not history:
        return "（暂无）"
    lines: List[str] = []
    for item in history[-limit:]:
        call = item.get("call", {})
        name = call.get("name", "未知工具")
        attempt = item.get("attempt", 0)
        success = item.get("success", False)
        status = "成功" if success else "失败"
        lines.append(f"- {name} 第 {attempt} 次 → {status}")
        if item.get("output"):
            lines.append("  输出：" + _truncate_text(item["output"], 200))
        if item.get("error"):
            lines.append("  错误：" + _truncate_text(item["error"], 200))
    return "\n".join(lines)


def _format_schema_parameters(schema: Dict[str, Any]) -> List[str]:
    if not schema:
        return ["  - 无参数"]
    props = schema.get("properties") or {}
    if not props:
        return ["  - 无参数"]
    required = set(schema.get("required") or [])
    lines: List[str] = []
    for field, meta in props.items():
        meta = meta or {}
        field_type = meta.get("type", "string")
        desc = (meta.get("description") or "").strip()
        req_label = "必填" if field in required else "可选"
        line = f"  - {field} ({field_type}，{req_label})"
        if desc:
            line += f"：{desc}"
        lines.append(line)
    return lines or ["  - 无参数"]


def _generate_example_args(schema: Dict[str, Any]) -> Dict[str, Any]:
    example: Dict[str, Any] = {}
    if not schema:
        return example
    props = schema.get("properties") or {}
    for field, meta in props.items():
        meta = meta or {}
        if "default" in meta:
            example[field] = meta["default"]
            continue
        enum_values = meta.get("enum") or []
        if enum_values:
            example[field] = enum_values[0]
            continue
        field_type = meta.get("type", "string")
        if field_type in ("number", "integer"):
            example[field] = 0
        elif field_type == "boolean":
            example[field] = True
        elif field_type == "array":
            example[field] = []
        elif field_type == "object":
            example[field] = {}
        else:
            description_hint = meta.get("description") or ""
            example[field] = description_hint or ""
    return example


def _format_tool_block(spec: WorkflowToolSpec) -> str:
    param_lines = _format_schema_parameters(spec.schema)
    example = json.dumps(spec.example_args, ensure_ascii=False) if spec.example_args else "{}"
    lines = [
        f"- 工具名：{spec.name}",
        f"  功能：{spec.description or '暂无描述'}",
        "  参数：",
        *param_lines,
        f"  示例：{example}",
    ]
    return "\n".join(lines)


def _build_workflow_tool_spec(tool_def: Dict[str, Any]) -> Optional[WorkflowToolSpec]:
    if not tool_def:
        return None
    name = tool_def.get("name")
    if not name:
        return None
    description = tool_def.get("description") or tool_def.get("summary") or ""
    schema = tool_def.get("inputSchema") or {}
    example_args = _generate_example_args(schema)

    def _executor(args: Dict[str, Any], attempt: int) -> Tuple[bool, Optional[str], Optional[str]]:
        try:
            resp = requests.post(
                f"http://127.0.0.1:5010/api/mcp/tools/{name}",
                json=args,
                timeout=120,
            )
            resp.raise_for_status()
            data = resp.json()
        except Exception as exc:
            util.log(1, f"调用工具 {name} 异常: {exc}")
            return False, None, str(exc)

        if data.get("success"):
            result = data.get("result")
            output = _normalize_tool_output(result)
            return True, output, None

        error_msg = data.get("error") or "未知错误"
        util.log(1, f"调用工具 {name} 失败: {error_msg}")
        return False, None, error_msg

    return WorkflowToolSpec(
        name=name,
        description=description,
        schema=schema,
        executor=_executor,
        example_args=example_args,
    )


def _build_builtin_tool_spec(
    name: str,
    description: str,
    schema: Dict[str, Any],
    executor: Callable[[Dict[str, Any], int], Tuple[bool, Optional[str], Optional[str]]],
) -> WorkflowToolSpec:
    return WorkflowToolSpec(
        name=name,
        description=description,
        schema=schema,
        executor=executor,
        example_args=_generate_example_args(schema),
    )


def _emotion_support_plan_executor(args: Dict[str, Any], attempt: int) -> Tuple[bool, Optional[str], Optional[str]]:
    valence = str(args.get("valence") or "neutral").lower()
    arousal = str(args.get("arousal") or "medium").lower()
    user_goal = str(args.get("user_goal") or "").strip()

    if valence == "negative" and arousal == "low":
        text = (
            "建议先做低刺激安抚：\n"
            "1. 先共情，不急着分析。\n"
            "2. 只给一个很小的下一步。\n"
            "3. 如果用户愿意，可顺带推荐舒缓音乐或帮他看近期安排。"
        )
    elif valence == "negative" and arousal == "high":
        text = (
            "建议先做稳定情绪：\n"
            "1. 先帮助用户慢下来，避免连续追问。\n"
            "2. 语言简短、确定、低刺激。\n"
            "3. 如用户愿意，可推荐舒缓音乐，或帮他查询/整理最近安排。"
        )
    elif valence == "positive":
        text = (
            "建议先共鸣再延展：\n"
            "1. 先接住用户的开心或期待。\n"
            "2. 再顺势问一个轻问题或给行动建议。\n"
            "3. 不要突然切成冷静说教模式。"
        )
    else:
        text = (
            "建议保持支持性回应：\n"
            "1. 先确认用户需求。\n"
            "2. 回答简洁清楚。\n"
            "3. 如果用户在安排事情，可考虑调用日程相关工具。"
        )

    if user_goal:
        text += f"\n当前用户目标线索：{user_goal}"
    return True, text, None


def _care_music_recommend_executor(args: Dict[str, Any], attempt: int) -> Tuple[bool, Optional[str], Optional[str]]:
    mood = str(args.get("mood") or "neutral").lower()
    energy = str(args.get("energy") or "medium").lower()
    language = str(args.get("language") or "中文").strip() or "中文"
    scenario = str(args.get("scenario") or "陪伴").strip() or "陪伴"

    if mood == "negative" and energy == "low":
        keywords = [
            f"{language} 钢琴 治愈 安静",
            f"{language} 轻音乐 放松 助眠",
            f"{language} acoustic 温柔 陪伴",
        ]
    elif mood == "negative" and energy == "high":
        keywords = [
            f"{language} 呼吸 放松 白噪音",
            f"{language} ambient calm down",
            f"{language} 冥想 放松 降躁",
        ]
    elif mood == "positive":
        keywords = [
            f"{language} 开心 元气 playlist",
            f"{language} 通勤 愉快 {scenario}",
            f"{language} upbeat feel good",
        ]
    else:
        keywords = [
            f"{language} 轻松 背景音乐",
            f"{language} 专注 陪伴 playlist",
            f"{language} 舒缓 日常",
        ]

    text = "可直接用于音乐平台搜索的关键词：\n" + "\n".join(f"- {item}" for item in keywords)
    return True, text, None


def _build_motivational_quote_key(username: str, model_id: str) -> str:
    return f"{username or 'User'}::{model_id or 'default'}"


def _get_motivational_quote_cooldown_status(username: Optional[str], model_id: Optional[str]) -> Tuple[bool, int]:
    key = _build_motivational_quote_key(username or "User", model_id or "default")
    now = time.time()
    with motivational_quote_lock:
        last_used = motivational_quote_cooldowns.get(key)
    if not last_used:
        return True, 0
    remaining = int(max(0, MOTIVATIONAL_QUOTE_MIN_INTERVAL_SECONDS - (now - last_used)))
    return remaining <= 0, remaining


def _touch_motivational_quote_cooldown(username: Optional[str], model_id: Optional[str]) -> None:
    key = _build_motivational_quote_key(username or "User", model_id or "default")
    with motivational_quote_lock:
        motivational_quote_cooldowns[key] = time.time()


def _infer_motivational_voice(role_profile: Dict[str, Any]) -> str:
    name = str(role_profile.get("name") or role_profile.get("first_name") or "").strip()
    position = str(role_profile.get("position") or "").strip()
    occupation = str(role_profile.get("occupation") or role_profile.get("job") or "").strip()
    additional = str(role_profile.get("additional") or "").strip()
    combined = f"{name} {position} {occupation} {additional}"

    historical_names = {
        "孔子", "老子", "庄子", "孟子", "荀子", "墨子", "韩非子",
        "李白", "杜甫", "苏轼", "李清照", "曹操", "刘备", "诸葛亮",
    }
    if name in historical_names:
        return "sage"
    if "陪伴" in combined or "温柔" in combined or "治愈" in combined:
        return "gentle"
    if "娱乐" in combined or "活泼" in combined or "外向" in combined:
        return "playful"
    if "教培" in combined or "老师" in combined or "教师" in combined or "顾问" in combined:
        return "mentor"
    if "客服" in combined or "助理" in combined or "专业" in combined:
        return "steady"
    return "natural"


def _render_role_based_quote(base_quote: str, role_profile: Dict[str, Any]) -> str:
    voice = _infer_motivational_voice(role_profile)
    name = str(role_profile.get("name") or role_profile.get("first_name") or "").strip()

    if voice == "sage":
        return f"{base_quote} 一步一步走，路会慢慢亮起来。"
    if voice == "gentle":
        return f"{base_quote} 先把眼前这一步处理好，我们再继续。"
    if voice == "playful":
        return f"{base_quote} 今天先赢回一点点状态，也算漂亮的一分。"
    if voice == "mentor":
        return f"{base_quote} 先把下一小步做出来，节奏就会回来。"
    if voice == "steady":
        return f"{base_quote} 先稳住当下，再处理后面的事。"
    if name:
        return f"{base_quote} {name}会陪你把这一步走过去。"
    return base_quote


def _motivational_quote_executor(args: Dict[str, Any], attempt: int) -> Tuple[bool, Optional[str], Optional[str]]:
    mood = str(args.get("mood") or "neutral").lower()
    energy = str(args.get("energy") or "medium").lower()
    scenario = str(args.get("scenario") or "日常陪伴").strip() or "日常陪伴"
    username = str(args.get("_username") or "User").strip() or "User"
    model_id = str(args.get("_model_id") or "default").strip() or "default"
    role_profile = args.get("_role_profile") if isinstance(args.get("_role_profile"), dict) else {}

    can_use, remaining = _get_motivational_quote_cooldown_status(username, model_id)
    if not can_use:
        wait_minutes = max(1, int((remaining + 59) // 60))
        text = (
            f"当前不适合再次输出励志名言。距离上次触发还需约 {wait_minutes} 分钟。\n"
            "使用建议：这次请直接延续自然陪伴、共情或简短支持，不要再额外插入名言。"
        )
        return True, text, None

    if mood == "negative" and energy == "low":
        quotes = [
            "你不用一下子走很远，先把眼前这一步走稳就很好。",
            "允许自己慢一点，但别轻易否定自己，你已经在往前了。",
            "今天状态不理想也正常，你能坚持到现在，本身就说明你很有力量。",
        ]
    elif mood == "negative" and energy == "high":
        quotes = [
            "先把呼吸放慢，稳住自己，很多难题都会在冷静后变得能处理。",
            "你不需要立刻赢下所有事情，先把最重要的一件做好就够了。",
            "情绪再满，也不代表你不行，你只是此刻需要一点缓冲和整理。",
        ]
    elif mood == "positive":
        quotes = [
            "你现在的这股劲头很珍贵，继续走，很多好结果都在路上。",
            "把今天这份状态记住，它会成为你下一次低谷时的底气。",
            "当你愿意认真生活的时候，生活通常也会慢慢给你回应。",
        ]
    else:
        quotes = [
            "别小看日复一日的普通努力，很多改变都是这样慢慢发生的。",
            "你可以不着急证明什么，先把自己的节奏照顾好。",
            "再普通的一天，也可以成为重新出发的起点。",
        ]

    index = max(0, (attempt - 1) % len(quotes))
    role_quote = _render_role_based_quote(quotes[index], role_profile)
    _touch_motivational_quote_cooldown(username, model_id)
    text = (
        f"适合当前{scenario}场景的一句励志名言：\n"
        f"{role_quote}\n"
        "使用建议：自然地夹在安抚、鼓励或收尾回应里，像角色本人会说的话，不要连续密集输出。"
    )
    return True, text, None


def _current_weather_executor(args: Dict[str, Any], attempt: int) -> Tuple[bool, Optional[str], Optional[str]]:
    location = str(args.get("location") or "").strip()
    if not location:
        return False, None, "location 不能为空"

    try:
        geo_resp = requests.get(
            "https://geocoding-api.open-meteo.com/v1/search",
            params={"name": location, "count": 1, "language": "zh", "format": "json"},
            timeout=15,
        )
        geo_resp.raise_for_status()
        geo_data = geo_resp.json() or {}
        results = geo_data.get("results") or []
        if not results:
            return False, None, f"没有找到地点：{location}"

        place = results[0]
        lat = place.get("latitude")
        lon = place.get("longitude")
        resolved_name = place.get("name") or location
        admin = place.get("admin1") or ""
        country = place.get("country") or ""

        weather_resp = requests.get(
            "https://api.open-meteo.com/v1/forecast",
            params={
                "latitude": lat,
                "longitude": lon,
                "current": "temperature_2m,relative_humidity_2m,apparent_temperature,precipitation,rain,showers,snowfall,weather_code,wind_speed_10m",
                "daily": "weather_code,temperature_2m_max,temperature_2m_min,precipitation_probability_max",
                "timezone": "auto",
                "forecast_days": 1,
            },
            timeout=15,
        )
        weather_resp.raise_for_status()
        weather_data = weather_resp.json() or {}
    except Exception as exc:
        return False, None, f"天气查询失败: {exc}"

    current = weather_data.get("current") or {}
    daily = weather_data.get("daily") or {}
    weather_code = current.get("weather_code")
    code_map = {
        0: "晴",
        1: "大部晴朗",
        2: "局部多云",
        3: "阴",
        45: "雾",
        48: "雾凇",
        51: "小毛毛雨",
        53: "毛毛雨",
        55: "强毛毛雨",
        61: "小雨",
        63: "中雨",
        65: "大雨",
        71: "小雪",
        73: "中雪",
        75: "大雪",
        80: "阵雨",
        81: "较强阵雨",
        82: "强阵雨",
        95: "雷暴",
        "current_weather": _build_builtin_tool_spec(
            "current_weather",
            "查询指定城市或地区的当前天气、体感温度、风速和今日温度范围。遇到天气、温度、下雨、气温问题时优先调用。",
            weather_schema,
            _current_weather_executor,
        ),
    }
    weather_text = code_map.get(weather_code, f"天气代码 {weather_code}")
    code_map.pop("current_weather", None)
    code_map = {key: value for key, value in code_map.items() if isinstance(key, int)}
    max_temp = (daily.get("temperature_2m_max") or [None])[0]
    min_temp = (daily.get("temperature_2m_min") or [None])[0]
    pop = (daily.get("precipitation_probability_max") or [None])[0]

    region = " ".join(part for part in [resolved_name, admin, country] if part).strip()
    lines = [
        f"{region} 当前天气：{weather_text}",
        f"当前温度 {current.get('temperature_2m')}°C，体感 {current.get('apparent_temperature')}°C",
        f"湿度 {current.get('relative_humidity_2m')}%，风速 {current.get('wind_speed_10m')} km/h",
    ]
    if max_temp is not None and min_temp is not None:
        lines.append(f"今天预计 {min_temp}°C - {max_temp}°C")
    if pop is not None:
        lines.append(f"降水概率约 {pop}%")
    return True, "\n".join(lines), None


def _get_builtin_tool_specs() -> Dict[str, WorkflowToolSpec]:
    support_schema = {
        "type": "object",
        "properties": {
            "valence": {"type": "string", "description": "negative / neutral / positive"},
            "arousal": {"type": "string", "description": "low / medium / high"},
            "user_goal": {"type": "string", "description": "用户当前目标或困扰"},
        },
    }
    music_schema = {
        "type": "object",
        "properties": {
            "mood": {"type": "string", "description": "negative / neutral / positive"},
            "energy": {"type": "string", "description": "low / medium / high"},
            "language": {"type": "string", "description": "推荐关键词使用的语言"},
            "scenario": {"type": "string", "description": "使用场景，如睡前、通勤、陪伴"},
        },
    }
    motivational_quote_schema = {
        "type": "object",
        "properties": {
            "mood": {"type": "string", "description": "negative / neutral / positive"},
            "energy": {"type": "string", "description": "low / medium / high"},
            "scenario": {"type": "string", "description": "使用场景，如安慰、鼓劲、受挫后、日常陪伴"},
        },
    }
    weather_schema = {
        "type": "object",
        "properties": {
            "location": {"type": "string", "description": "要查询天气的城市或地区，例如 济南、上海、北京"},
        },
        "required": ["location"],
    }
    return {
        "emotion_support_plan": _build_builtin_tool_spec(
            "emotion_support_plan",
            "根据当前情绪状态给出一份简短的关怀和回复策略，适合在安慰、稳定情绪、共鸣时调用。",
            support_schema,
            _emotion_support_plan_executor,
        ),
        "care_music_recommend": _build_builtin_tool_spec(
            "care_music_recommend",
            "给出适合当前情绪状态的音乐搜索关键词，用于情绪安抚或陪伴场景。",
            music_schema,
            _care_music_recommend_executor,
        ),
        "motivational_quote": _build_builtin_tool_spec(
            "motivational_quote",
            "在合适的陪伴时机给出一句简短自然的励志名言，适合低落、受挫、需要打气或日常鼓励场景。",
            motivational_quote_schema,
            _motivational_quote_executor,
        ),
        "current_weather": _build_builtin_tool_spec(
            "current_weather",
            "查询指定城市或地区的当前天气、体感温度、风速和今日温度范围。遇到天气、温度、下雨、气温问题时优先调用。",
            weather_schema,
            _current_weather_executor,
        ),
    }


def _build_emotion_policy_suffix(observation: Optional[str]) -> str:
    text = observation or ""
    if "emotion_state:" not in text:
        return (
            "\n【情绪决策】如果额外观察里包含用户情绪线索，请先根据情绪调整回复方式："
            "低落时先确认感受并给一个明确下一步，焦虑时先稳定节奏，开心时先共鸣，再决定是否给建议或调用工具。"
            "避免幼态表达与模板安抚句，禁止使用“你好呀/慢慢来/不急的/没关系啦”等话术。"
        )

    return (
        "\n【情绪决策】你已获得结构化 emotion_state。请严格按以下顺序处理："
        "先接住情绪，再回答问题；negative/low 时优先确认感受并给具体建议，negative/high 时优先稳定情绪；"
        "positive 时先共鸣再延展。若 tool_suggestions 包含 care_music_recommend、motivational_quote 或日程工具，且对用户有帮助，可主动调用；"
        "其中 motivational_quote 只适合偶尔点到为止，不要高频重复。"
        "禁止输出幼态称呼与模板安抚句（如“你好呀”“慢慢来”“不急的”“没关系啦”）。"
    )


def _format_tools_for_prompt(tool_specs: Dict[str, WorkflowToolSpec]) -> str:
    if not tool_specs:
        return "（暂无可用工具）"
    return "\n".join(_format_tool_block(spec) for spec in tool_specs.values())


def _build_planner_messages(state: AgentState) -> List[SystemMessage | HumanMessage]:
    context = state.get("context", {}) or {}
    system_prompt = context.get("system_prompt", "")
    request = state.get("request", "")
    tool_specs = context.get("tool_registry", {}) or {}
    planner_preview = state.get("planner_preview")
    conversation = state.get("messages", []) or []
    history = state.get("tool_results", []) or []
    memory_context = context.get("memory_context", "")
    knowledge_context = context.get("knowledge_context", "")
    observation = context.get("observation", "")

    convo_text = "\n".join(f"{msg['role']}: {msg['content']}" for msg in conversation) or "（暂无对话）"
    history_text = _truncate_history(history)
    tools_text = _format_tools_for_prompt(tool_specs)
    preview_section = f"\n（规划器预览：{planner_preview}）" if planner_preview else ""
    username = context.get("username")
    model_id = context.get("model_id")
    can_use_quote, quote_remaining = _get_motivational_quote_cooldown_status(username, model_id)
    quote_hint = (
        "励志名言工具当前可用；仅在明显需要鼓励、打气、安抚时再调用。"
        if can_use_quote
        else f"励志名言工具仍在冷却中，约 {max(1, (quote_remaining + 59) // 60)} 分钟后再考虑调用。"
    )

    user_block = textwrap.dedent(
        f"""
你是一名助理，请根据请求决定是否需要调用工具。

当前请求：
{request}

系统设定：
{system_prompt}

额外观察：
{observation or '（无补充）'}

相关记忆：
{memory_context or '（无相关记忆）'}

相关知识：
{knowledge_context or '（无相关知识）'}

对话及工具记录：
{convo_text}

可用工具：
{tools_text}

补充约束：
{quote_hint}

历史工具执行：
{history_text}{preview_section}

请返回 JSON，格式如下：
- 若需要调用工具：
    {{"action": "tool", "tool": "工具名", "args": {{...}}}}
- 若直接回复：
    {{"action": "finish_text"}}
        """
    ).strip()

    return [
        SystemMessage(content="你负责规划下一步行动，请严格输出合法 JSON。"),
        HumanMessage(content=user_block),
    ]


def _build_final_messages(state: AgentState) -> List[SystemMessage | HumanMessage]:
    context = state.get("context", {}) or {}
    system_prompt = context.get("system_prompt", "")
    knowledge_context = context.get("knowledge_context", "")
    memory_context = context.get("memory_context", "")
    observation = context.get("observation", "")
    conversation = state.get("messages", []) or []
    planner_preview = state.get("planner_preview")
    conversation_block = "\n".join(f"{msg['role']}: {msg['content']}" for msg in conversation) or "（暂无对话）"
    history_text = _truncate_history(state.get("tool_results", []))
    preview_section = f"\n（规划器建议：{planner_preview}）" if planner_preview else ""

    user_block = textwrap.dedent(
        f"""
系统设定：
{system_prompt}

相关记忆：
{memory_context or '（无相关记忆）'}

相关知识：
{knowledge_context or '（无相关知识）'}

其他观察：
{observation or '（无补充）'}

对话及工具记录：
{conversation_block}

工具执行摘要：
{history_text}{preview_section}

请结合以上信息，为用户生成自然语言答复，保持人设与语气一致。
        """
    ).strip()

    return [
        SystemMessage(content="你是最终回复的口播助手，请用中文自然表达。请保持回复简短，控制在2-4句话或约100字以内，方便阅读。"),
        HumanMessage(content=user_block),
    ]


def _call_planner_llm(state: AgentState) -> Dict[str, Any]:
    response = llm.invoke(_build_planner_messages(state))
    content = getattr(response, "content", None)
    if not isinstance(content, str):
        raise RuntimeError("规划器返回内容异常，未获得字符串。")
    trimmed = content.strip()
    try:
        decision = json.loads(trimmed)
    except json.JSONDecodeError:
        decision = {"action": "finish", "message": trimmed}
    decision.setdefault("_raw", trimmed)
    return decision


def _plan_next_action(state: AgentState) -> AgentState:
    context = state.get("context", {}) or {}
    audit_log = list(state.get("audit_log", []))
    history = state.get("tool_results", []) or []
    max_steps = state.get("max_steps", 12)
    if len(history) >= max_steps:
        audit_log.append("规划器：超过最大步数，终止流程。")
        return {
            "status": "failed",
            "audit_log": audit_log,
            "error": "工具调用步数超限",
            "context": context,
        }

    decision = _call_planner_llm(state)
    audit_log.append(f"规划器：决策 -> {decision.get('_raw', decision)}")

    action = decision.get("action")
    if action == "tool":
        tool_name = decision.get("tool")
        tool_registry: Dict[str, WorkflowToolSpec] = context.get("tool_registry", {})
        if tool_name not in tool_registry:
            audit_log.append(f"规划器：未知工具 {tool_name}")
            return {
                "status": "failed",
                "audit_log": audit_log,
                "error": f"未知工具 {tool_name}",
                "context": context,
            }
        args = decision.get("args") or {}

        if history:
            last_entry = history[-1]
            last_call = last_entry.get("call", {}) or {}
            if (
                last_entry.get("success")
                and last_call.get("name") == tool_name
                and (last_call.get("args") or {}) == args
                and last_entry.get("output")
            ):
                recent_attempts = sum(
                    1
                    for item in reversed(history)
                    if item.get("call", {}).get("name") == tool_name
                )
                if recent_attempts >= 1:
                    audit_log.append(
                        "规划器：检测到工具重复调用，使用最新结果产出最终回复。"
                    )
                    final_messages = _build_final_messages(state)
                    preview = last_entry.get("output")
                    return {
                        "status": "completed",
                        "planner_preview": preview,
                        "final_response": None,
                        "final_messages": final_messages,
                        "audit_log": audit_log,
                        "context": context,
                    }
        return {
            "next_action": {"name": tool_name, "args": args},
            "status": "needs_tool",
            "audit_log": audit_log,
            "context": context,
        }

    if action in {"finish", "finish_text"}:
        preview = decision.get("message")
        final_messages = _build_final_messages(state)
        audit_log.append("规划器：任务完成，准备输出最终回复。")
        return {
            "status": "completed",
            "planner_preview": preview,
            "final_response": preview if action == "finish" else None,
            "final_messages": final_messages,
            "audit_log": audit_log,
            "context": context,
        }

    raise RuntimeError(f"未知的规划器决策: {decision}")


def _execute_tool(state: AgentState) -> AgentState:
    context = dict(state.get("context", {}) or {})
    action = state.get("next_action")
    if not action:
        return {
            "status": "failed",
            "error": "缺少要执行的工具指令",
            "context": context,
        }

    history = list(state.get("tool_results", []) or [])
    audit_log = list(state.get("audit_log", []) or [])
    conversation = list(state.get("messages", []) or [])

    name = action.get("name")
    args = dict(action.get("args", {}) or {})
    tool_registry: Dict[str, WorkflowToolSpec] = context.get("tool_registry", {})
    spec = tool_registry.get(name)
    if not spec:
        return {
            "status": "failed",
            "error": f"未知工具 {name}",
            "context": context,
        }

    attempts = sum(1 for item in history if item.get("call", {}).get("name") == name)
    if name == "motivational_quote":
        args.setdefault("_username", context.get("username"))
        args.setdefault("_model_id", context.get("model_id"))
        args.setdefault("_role_profile", context.get("role_profile"))
    success, output, error = spec.executor(args, attempts)
    result: ToolResult = {
        "call": {"name": name, "args": args},
        "success": success,
        "output": output,
        "error": error,
        "attempt": attempts + 1,
    }
    history.append(result)
    audit_log.append(f"执行器：{name} 第 {result['attempt']} 次 -> {'成功' if success else '失败'}")

    message_lines = [
        f"[TOOL] {name} {'成功' if success else '失败'}。",
    ]
    if output:
        message_lines.append(f"[TOOL] 输出：{_truncate_text(output, 200)}")
    if error:
        message_lines.append(f"[TOOL] 错误：{_truncate_text(error, 200)}")
    conversation.append({"role": "assistant", "content": "\n".join(message_lines)})

    return {
        "tool_results": history,
        "messages": conversation,
        "next_action": None,
        "audit_log": audit_log,
        "status": "planning",
        "error": error if not success else None,
        "context": context,
    }


def _route_decision(state: AgentState) -> str:
    return "call_tool" if state.get("status") == "needs_tool" else "end"


def _build_workflow_app() -> StateGraph:
    graph = StateGraph(AgentState)
    graph.add_node("plan_next", _plan_next_action)
    graph.add_node("call_tool", _execute_tool)
    graph.add_edge(START, "plan_next")
    graph.add_conditional_edges(
        "plan_next",
        _route_decision,
        {
            "call_tool": "call_tool",
            "end": END,
        },
    )
    graph.add_edge("call_tool", "plan_next")
    return graph.compile()


_WORKFLOW_APP = _build_workflow_app()

def get_user_memory_dir(username=None, model_id=None):
    """根据配置决定是否按用户名和模型ID隔离记忆目录"""
    if username is None:
        username = current_username
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    mem_base = os.path.join(base_dir, "memory")
    try:
        cfg.load_config()
        isolate = cfg.config["memory"]["isolate_by_user"]
    except Exception:
        isolate = False
    
    # 如果启用了用户隔离
    if isolate and username:
        user_dir = os.path.join(mem_base, str(username))
        # 如果提供了模型ID，进一步按模型隔离
        if model_id:
            return os.path.join(user_dir, str(model_id))
        return user_dir
    
    # 如果没有启用用户隔离，但提供了模型ID，仍然按模型隔离
    if model_id:
        return os.path.join(mem_base, str(model_id))
    
    return mem_base

def get_current_time_step(username=None):
    """
    获取当前时间作为time_step
    
    返回:
        int: 当前时间步，从0开始，非真实时间
    """
    global agents
    try:
        # 按用户名选择对应agent，若未指定则退回全局agent
        ag = agents.get(username) if username else None
        if ag and ag.memory_stream and ag.memory_stream.seq_nodes:
            # 如果有记忆节点，则使用最后一个节点的created属性加1
            return int(ag.memory_stream.seq_nodes[-1].created) + 1
        else:
            # 如果没有记忆节点或agent未初始化，则使用0
            return 0
    except Exception as e:
        util.log(1, f"获取time_step时出错: {str(e)}，使用0代替")
        return 0

# 新增：本地知识库相关函数
def read_doc_file(file_path):
    """
    读取doc文件内容
    
    参数:
        file_path: doc文件路径
        
    返回:
        str: 文档内容
    """
    try:
        # 方法1: 使用 win32com.client（Windows系统，推荐用于.doc文件）
        if WIN32COM_AVAILABLE:
            word = None
            doc = None
            try:
                import pythoncom
                pythoncom.CoInitialize()  # 初始化COM组件
                
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                doc = word.Documents.Open(file_path)
                content = doc.Content.Text
                
                # 先保存内容，再尝试关闭
                if content and content.strip():
                    try:
                        doc.Close()
                        word.Quit()
                    except Exception as close_e:
                        util.log(1, f"关闭Word应用程序时出错: {str(close_e)}，但内容已成功提取")
                    
                    try:
                        pythoncom.CoUninitialize()  # 清理COM组件
                    except:
                        pass
                    
                    return content.strip()
                
            except Exception as e:
                util.log(1, f"使用 win32com 读取 .doc 文件失败: {str(e)}")
            finally:
                # 确保资源被释放
                try:
                    if doc:
                        doc.Close()
                except:
                    pass
                try:
                    if word:
                        word.Quit()
                except:
                    pass
                try:
                    pythoncom.CoUninitialize()
                except:
                    pass
        
        # 方法2: 简单的二进制文本提取（备选方案）
        try:
            with open(file_path, 'rb') as f:
                raw_data = f.read()
                # 尝试提取可打印的文本
                text_parts = []
                current_text = ""
                
                for byte in raw_data:
                    char = chr(byte) if 32 <= byte <= 126 or byte in [9, 10, 13] else None
                    if char:
                        current_text += char
                    else:
                        if len(current_text) > 3:  # 只保留长度大于3的文本片段
                            text_parts.append(current_text.strip())
                        current_text = ""
                
                if len(current_text) > 3:
                    text_parts.append(current_text.strip())
                
                # 过滤和清理文本
                filtered_parts = []
                for part in text_parts:
                    # 移除过多的重复字符和无意义的片段
                    if (len(part) > 5 and 
                        not part.startswith('Microsoft') and 
                        not all(c in '0123456789-_.' for c in part) and
                        len(set(part)) > 3):  # 字符种类要多样
                        filtered_parts.append(part)
                
                if filtered_parts:
                    return '\n'.join(filtered_parts)
                    
        except Exception as e:
            util.log(1, f"使用二进制方法读取 .doc 文件失败: {str(e)}")
        
        util.log(1, f"无法读取 .doc 文件 {file_path}，建议转换为 .docx 格式")
        return ""
        
    except Exception as e:
        util.log(1, f"读取doc文件 {file_path} 时出错: {str(e)}")
        return ""

def read_docx_file(file_path):
    """
    读取docx文件内容
    
    参数:
        file_path: docx文件路径
        
    返回:
        str: 文档内容
    """
    try:
        doc = docx.Document(file_path)
        content = []
        
        for element in doc.element.body:
            if isinstance(element, CT_P):
                paragraph = Paragraph(element, doc)
                if paragraph.text.strip():
                    content.append(paragraph.text.strip())
            elif isinstance(element, CT_Tbl):
                table = Table(element, doc)
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        content.append(" | ".join(row_text))
        
        return "\n".join(content)
    except Exception as e:
        util.log(1, f"读取docx文件 {file_path} 时出错: {str(e)}")
        return ""
    
def read_pptx_file(file_path):
    """
    读取pptx文件内容
    
    参数:
        file_path: pptx文件路径
        
    返回:
        str: 演示文稿内容
    """
    if not PPTX_AVAILABLE:
        util.log(1, "python-pptx 库未安装，无法读取 PowerPoint 文件")
        return ""
        
    try:
        prs = Presentation(file_path)
        content = []
        
        for i, slide in enumerate(prs.slides):
            slide_content = [f"第{i+1}页："]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
                    
            if len(slide_content) > 1:  # 有内容才添加
                content.append("\n".join(slide_content))
        
        return "\n\n".join(content)
    except Exception as e:
        util.log(1, f"读取pptx文件 {file_path} 时出错: {str(e)}")
        return ""

def load_local_knowledge_base():
    """
    加载本地知识库内容
    
    返回:
        dict: 文件名到内容的映射
    """
    knowledge_base = {}
    
    # 获取llm/data目录路径
    current_dir = os.path.dirname(os.path.abspath(__file__))
    data_dir = os.path.join(current_dir, "data")
    
    if not os.path.exists(data_dir):
        util.log(1, f"知识库目录不存在: {data_dir}")
        return knowledge_base
    
    # 遍历data目录中的文件
    for file_path in Path(data_dir).iterdir():
        if not file_path.is_file():
            continue
            
        file_name = file_path.name
        file_extension = file_path.suffix.lower()
        
        try:
            if file_extension == '.docx':
                content = read_docx_file(str(file_path))
            elif file_extension == '.doc':
                content = read_doc_file(str(file_path))
            elif file_extension == '.pptx':
                content = read_pptx_file(str(file_path))
            else:
                # 尝试作为文本文件读取
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                except UnicodeDecodeError:
                    try:
                        with open(file_path, 'r', encoding='gbk') as f:
                            content = f.read()
                    except UnicodeDecodeError:
                        util.log(1, f"无法解码文件: {file_name}")
                        continue
            
            if content.strip():
                knowledge_base[file_name] = content
                util.log(1, f"成功加载知识库文件: {file_name} ({len(content)} 字符)")
            
        except Exception as e:
            util.log(1, f"加载知识库文件 {file_name} 时出错: {str(e)}")
    
    return knowledge_base

def search_knowledge_base(query, knowledge_base, max_results=3):
    """
    在知识库中搜索相关内容
    
    参数:
        query: 查询内容
        knowledge_base: 知识库字典
        max_results: 最大返回结果数
        
    返回:
        list: 相关内容列表
    """
    if not knowledge_base:
        return []
    
    results = []
    query_lower = query.lower()
    
    # 搜索关键词
    query_keywords = re.findall(r'\w+', query_lower)
    
    for file_name, content in knowledge_base.items():
        content_lower = content.lower()
        
        # 计算匹配度
        score = 0
        matched_sentences = []
        
        # 按句子分割内容
        sentences = re.split(r'[。！？\n]', content)
        
        for sentence in sentences:
            if not sentence.strip():
                continue
                
            sentence_lower = sentence.lower()
            sentence_score = 0
            
            # 计算关键词匹配度
            for keyword in query_keywords:
                if keyword in sentence_lower:
                    sentence_score += 1
            
            # 如果句子有匹配，记录
            if sentence_score > 0:
                matched_sentences.append((sentence.strip(), sentence_score))
                score += sentence_score
        
        # 如果有匹配的内容
        if score > 0:
            # 按匹配度排序句子
            matched_sentences.sort(key=lambda x: x[1], reverse=True)
            
            # 取前几个最相关的句子
            relevant_sentences = [sent[0] for sent in matched_sentences[:5] if sent[0]]
            
            if relevant_sentences:
                results.append({
                    'file_name': file_name,
                    'score': score,
                    'content': '\n'.join(relevant_sentences)
                })
    
    # 按匹配度排序
    results.sort(key=lambda x: x['score'], reverse=True)
    
    return results[:max_results]

# 全局知识库缓存
_knowledge_base_cache = None
_knowledge_base_load_time = None
_knowledge_base_file_times = {}  # 存储文件的最后修改时间

def check_knowledge_base_changes():
    """
    检查知识库文件是否有变化
    
    返回:
        bool: 如果有文件变化返回True，否则返回False
    """
    global _knowledge_base_file_times
    
    # 获取llm/data目录路径
    current_dir = os.path.dirname(os.path.abspath(__file__))
    data_dir = os.path.join(current_dir, "data")
    
    if not os.path.exists(data_dir):
        return False
    
    current_file_times = {}
    
    # 遍历data目录中的文件
    for file_path in Path(data_dir).iterdir():
        if not file_path.is_file():
            continue
        
        file_name = file_path.name
        file_extension = file_path.suffix.lower()
        
        # 只检查支持的文件格式
        if file_extension in ['.docx', '.doc', '.pptx', '.txt'] or file_extension == '':
            try:
                mtime = os.path.getmtime(str(file_path))
                current_file_times[file_name] = mtime
            except OSError:
                continue
    
    # 检查是否有变化
    if not _knowledge_base_file_times:
        # 第一次检查，保存文件时间
        _knowledge_base_file_times = current_file_times
        return True
    
    # 比较文件时间
    if set(current_file_times.keys()) != set(_knowledge_base_file_times.keys()):
        # 文件数量发生变化
        _knowledge_base_file_times = current_file_times
        return True
    
    for file_name, mtime in current_file_times.items():
        if file_name not in _knowledge_base_file_times or _knowledge_base_file_times[file_name] != mtime:
            # 文件被修改
            _knowledge_base_file_times = current_file_times
            return True
    
    return False

def init_knowledge_base():
    """
    初始化知识库，在系统启动时调用
    """
    global _knowledge_base_cache, _knowledge_base_load_time
    
    util.log(1, "初始化本地知识库...")
    _knowledge_base_cache = load_local_knowledge_base()
    _knowledge_base_load_time = time.time()
    
    # 初始化文件修改时间跟踪
    check_knowledge_base_changes()
    
    util.log(1, f"知识库初始化完成，共 {len(_knowledge_base_cache)} 个文件")

def get_knowledge_base():
    """
    获取知识库，使用缓存机制
    
    返回:
        dict: 知识库内容
    """
    global _knowledge_base_cache, _knowledge_base_load_time
    
    # 如果缓存为空，先初始化
    if _knowledge_base_cache is None:
        init_knowledge_base()
        return _knowledge_base_cache
    
    # 检查文件是否有变化
    if check_knowledge_base_changes():
        util.log(1, "检测到知识库文件变化，正在重新加载...")
        _knowledge_base_cache = load_local_knowledge_base()
        _knowledge_base_load_time = time.time()
        util.log(1, f"知识库重新加载完成，共 {len(_knowledge_base_cache)} 个文件")
    
    return _knowledge_base_cache


# 定时保存记忆的线程
def memory_scheduler_thread():
    """
    定时任务线程，运行schedule调度器
    """
    while True:
        schedule.run_pending()
        time.sleep(60)  # 每分钟检查一次是否有定时任务需要执行

# 初始化定时保存记忆的任务
def init_memory_scheduler():
    """
    初始化定时保存记忆的任务
    """
    global agents
    
    # 确保agent已经创建
    if not agents:
        util.log(1, '创建代理实例...')
        create_agent()
    
    # 设置每天0点保存记忆
    schedule.every().day.at("00:00").do(save_agent_memory)
    
    # 设置每天晚上11点执行反思
    schedule.every().day.at("23:00").do(perform_daily_reflection)
    
    # 启动定时任务线程
    scheduler_thread = MyThread(target=memory_scheduler_thread)
    scheduler_thread.start()
    
    util.log(1, '定时任务已启动：每天0点保存记忆，每天23点执行反思')

def check_memory_files(username=None, model_id=None):
    """
    检查memory目录及其必要文件是否存在
    
    参数:
        username: 用户名
        model_id: 模型ID
    
    返回:
        memory_dir: memory目录路径
        is_complete: 是否已经存在完整的memory目录结构
    """
    
    # 根据配置与用户名、模型ID获取memory目录路径
    memory_dir = get_user_memory_dir(username, model_id)

    # 检查memory目录是否存在，不存在则创建
    if not os.path.exists(memory_dir):
        os.makedirs(memory_dir)
        util.log(1, f"创建memory目录: {memory_dir}")
    
    # 删除.memory_cleared标记文件（如果存在）
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    mem_base = os.path.join(base_dir, "memory")
    memory_cleared_flag_file = os.path.join(mem_base, ".memory_cleared")
    if os.path.exists(memory_cleared_flag_file):
        try:
            os.remove(memory_cleared_flag_file)
            util.log(1, f"清除删除记忆标记文件: {memory_cleared_flag_file}")
            # 重置记忆清除标记
            global memory_cleared
            memory_cleared = False
        except Exception as e:
            util.log(1, f"清除删除记忆标记文件时出错: {str(e)}")
    
    # 检查meta.json是否存在
    meta_file = os.path.join(memory_dir, "meta.json")
    is_complete = os.path.exists(meta_file)
    
    # 检查memory_stream目录是否存在，不存在则创建
    memory_stream_dir = os.path.join(memory_dir, "memory_stream")
    if not os.path.exists(memory_stream_dir):
        os.makedirs(memory_stream_dir)
        util.log(1, f"创建memory_stream目录: {memory_stream_dir}")
    
    # 检查必要的文件是否存在
    embeddings_path = os.path.join(memory_stream_dir, "embeddings.json")
    nodes_path = os.path.join(memory_stream_dir, "nodes.json")
    
    # 检查文件是否存在且不为空
    is_complete = (os.path.exists(embeddings_path) and os.path.getsize(embeddings_path) > 2 and
                  os.path.exists(nodes_path) and os.path.getsize(nodes_path) > 2)
    
    # 如果文件不存在，创建空的JSON文件
    if not os.path.exists(embeddings_path):
        with open(embeddings_path, 'w', encoding='utf-8') as f:
            f.write('{}')
    
    if not os.path.exists(nodes_path):
        with open(nodes_path, 'w', encoding='utf-8') as f:
            f.write('[]')
    
    return memory_dir, is_complete

def create_agent(username=None, model_id=None):
    """
    创建一个GenerativeAgent实例
    
    参数:
        username: 用户名
        model_id: 模型ID（可选，如果提供则从数据库加载模型属性）
    
    返回:
        agent: GenerativeAgent对象
    """
    global agents
    
    if username is None:
        username = "User"
    
    # 构建agent的缓存key（包含用户名和模型ID）
    agent_key = f"{username}_{model_id}" if model_id else username
    
    # 创建/复用代理
    with agent_lock:
        if agent_key in agents:
            return agents[agent_key]
        
        memory_dir, is_exist = check_memory_files(username, model_id)
        agent = GenerativeAgent(memory_dir)
        
        # 检查是否有scratch属性，如果没有则添加
        if not hasattr(agent, 'scratch'):
            agent.scratch = {}
        
        # 优先从数据库加载模型属性
        if model_id:
            try:
                from core import model_attribute_service
                attributes = model_attribute_service.get_model_attributes(model_id)
                if attributes:
                    # 应用模型属性到agent
                    model_attribute_service.apply_model_attributes(agent, model_id, attributes)
                    util.log(1, f"[创建Agent] 从数据库加载模型属性: {model_id}")
                else:
                    # 如果加载失败，回退到config.json
                    util.log(1, f"[创建Agent] 无法加载模型属性，使用config.json")
                    _load_attributes_from_config(agent)
            except Exception as e:
                util.log(1, f"[创建Agent] 加载模型属性失败: {e}，使用config.json")
                _load_attributes_from_config(agent)
        else:
            # 如果没有提供model_id，从config.json加载
            _load_attributes_from_config(agent)
        
        # 如果memory目录存在且不为空，则加载之前保存的记忆（不包括scratch数据）
        if is_exist:
            load_agent_memory(agent, username, model_id)
        
        # 缓存到字典
        agents[agent_key] = agent
    
    return agent


def _load_attributes_from_config(agent):
    """从config.json加载属性到agent"""
    try:
        cfg.load_config()
        scratch_data = {
            "first_name": cfg.config.get("attribute", {}).get("name", "Fay"),
            "last_name": "",
            "age": cfg.config.get("attribute", {}).get("age", "成年"),
            "sex": cfg.config.get("attribute", {}).get("gender", "女"),
            "additional": cfg.config.get("attribute", {}).get("additional", "友好、乐于助人"),
            "birthplace": cfg.config.get("attribute", {}).get("birth", ""),
            "position": cfg.config.get("attribute", {}).get("position", ""),
            "zodiac": cfg.config.get("attribute", {}).get("zodiac", ""),
            "constellation": cfg.config.get("attribute", {}).get("constellation", ""),
            "contact": cfg.config.get("attribute", {}).get("contact", ""),
            "voice": cfg.config.get("attribute", {}).get("voice", "zhifeng_emo"),
            "goal": cfg.config.get("attribute", {}).get("goal", ""),
            "occupation": cfg.config.get("attribute", {}).get("job", "助手"),
            "current_time": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        }
        agent.scratch = scratch_data
    except Exception as e:
        util.log(1, f"[创建Agent] 从config.json加载属性失败: {e}")
        # 使用默认值
        agent.scratch = {
            "first_name": "Fay",
            "last_name": "",
            "age": "成年",
            "sex": "女",
            "additional": "友好、乐于助人",
            "birthplace": "",
            "position": "",
            "zodiac": "",
            "constellation": "",
            "contact": "",
            "voice": "zhifeng_emo",
            "goal": "",
            "occupation": "助手",
            "current_time": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        }

def load_agent_memory(agent, username=None, model_id=None):
    """
    从文件加载代理的记忆
    
    参数:
        agent: GenerativeAgent对象
        username: 用户名
        model_id: 模型ID
    """
    try:
        # 获取memory目录路径（按需隔离）
        memory_dir = get_user_memory_dir(username, model_id)
        memory_stream_dir = os.path.join(memory_dir, "memory_stream")
        
        # 加载nodes.json
        nodes_path = os.path.join(memory_stream_dir, "nodes.json")
        if os.path.exists(nodes_path) and os.path.getsize(nodes_path) > 2:  # 文件存在且不为空
            with open(nodes_path, 'r', encoding='utf-8') as f:
                nodes_data = json.load(f)
                
                # 清空当前的seq_nodes
                agent.memory_stream.seq_nodes = []
                agent.memory_stream.id_to_node = {}
                
                # 重新创建节点
                for node_dict in nodes_data:
                    new_node = ConceptNode(node_dict)
                    agent.memory_stream.seq_nodes.append(new_node)
                    agent.memory_stream.id_to_node[new_node.node_id] = new_node
        
        # 加载embeddings.json
        embeddings_path = os.path.join(memory_stream_dir, "embeddings.json")
        if os.path.exists(embeddings_path) and os.path.getsize(embeddings_path) > 2:  # 文件存在且不为空
            with open(embeddings_path, 'r', encoding='utf-8') as f:
                embeddings_data = json.load(f)
                agent.memory_stream.embeddings = embeddings_data
        
        util.log(1, f"已加载代理记忆")
    except Exception as e:
        util.log(1, f"加载代理记忆失败: {str(e)}")

# 记忆对话内容的线程函数
def remember_conversation_thread(username, content, response_text, model_id=None):
    """
    在单独线程中记录对话内容到代理记忆
    
    参数:
        username: 用户名
        content: 用户问题内容
        response_text: 代理回答内容
        model_id: 模型ID（可选）
    """
    global agents
    try:
        with agent_lock:
            # 构建agent的缓存key
            agent_key = f"{username}_{model_id}" if model_id else username
            ag = agents.get(agent_key)
            if ag is None:
                return
            time_step = get_current_time_step(username)
            name = "主人" if username == "User" else username
            # 记录对话内容
            memory_content = f"在对话中，我回答了{name}的问题：{content}\n，我的回答是：{response_text}"
            ag.remember(memory_content, time_step)
    except Exception as e:
        util.log(1, f"记忆对话内容出错: {str(e)}")

def question(content, username, observation=None, pure_mode=False, system_override=None):
    """处理用户提问并返回回复。system_override: 前端传入的系统指令，会追加到系统提示末尾。"""
    global agents, current_username
    current_username = username
    observation = "" if observation is None else str(observation)
    full_response_text = ""
    accumulated_text = ""
    default_punctuations = [",", ".", "!", "?", "\n", "\uFF0C", "\u3002", "\uFF01", "\uFF1F"]
    is_first_sentence = True

    from core import stream_manager
    sm = stream_manager.new_instance()
    conversation_id = sm.get_conversation_id(username)

    # 如果是纯模型模式，跳过角色设定和记忆加载
    if pure_mode:
        util.log(1, f"[纯模型模式] 用户 {username} 启用纯模型对话模式")
        
        # 简单的系统提示，不包含角色设定
        system_prompt = "你是 SoulLink 的心理陪伴型数字人，请用自然、温和、具体的中文回应用户。"
        
        # 不加载记忆和知识库
        memory_context = ""
        knowledge_context = ""
        
        # 简单的对话历史（按模型筛选）
        try:
            # 获取用户当前选择的模型ID
            model_id = None
            try:
                from core import member_db
                member_db_instance = member_db.new_instance()
                model_id = member_db_instance.get_current_model(username)
            except Exception:
                pass
            history_records = content_db.new_instance().get_recent_messages_by_user(username=username, limit=10, model_id=model_id)
        except Exception as exc:
            util.log(1, f"加载历史消息失败: {exc}")
            history_records = []

        messages_buffer: List[ConversationMessage] = []

        def append_to_buffer(role: str, text_value: str) -> None:
            if not text_value:
                return
            messages_buffer.append({"role": role, "content": text_value})
            if len(messages_buffer) > 20:  # 纯模型模式使用更少的历史记录
                del messages_buffer[:-20]

        for msg_type, msg_text in history_records:
            role = 'assistant'
            if msg_type and msg_type.lower() in ('member', 'user'):
                role = 'user'
            append_to_buffer(role, msg_text)

        if (
            not messages_buffer
            or messages_buffer[-1]['role'] != 'user'
            or messages_buffer[-1]['content'] != content
        ):
            append_to_buffer('user', content)

        if SOULLINK_THERAPY_SUFFIX not in system_prompt:
            system_prompt = f"{system_prompt}\n\n{SOULLINK_THERAPY_SUFFIX}"
        if SOULLINK_THERAPY_SUFFIX not in system_prompt:
            system_prompt = f"{system_prompt}\n\n{SOULLINK_THERAPY_SUFFIX}"
        system_prompt = f"{system_prompt}{_build_emotion_policy_suffix(observation)}"
        messages = [SystemMessage(content=system_prompt), HumanMessage(content=content)]
        
        if system_override:
            system_prompt = f"{system_prompt}\n\n【补充】{system_override}"
        # 纯模型模式不使用工具
        tool_registry: Dict[str, WorkflowToolSpec] = _get_builtin_tool_specs()
        
    else:
        # 原有的角色模式逻辑
        # 获取用户当前选择的模型ID
        model_id = None
        try:
            from core import member_db
            member_db_instance = member_db.new_instance()
            model_id = member_db_instance.get_current_model(username)
        except Exception as e:
            util.log(1, f"[对话] 获取用户模型失败: {e}，使用默认配置")
        
        agent = create_agent(username, model_id)

        agent_desc = {
            "first_name": agent.scratch.get("first_name", "Fay"),
            "last_name": agent.scratch.get("last_name", ""),
            "age": agent.scratch.get("age", "成年"),
            "sex": agent.scratch.get("sex", "女"),
            "additional": agent.scratch.get("additional", "友好、乐于助人"),
            "birthplace": agent.scratch.get("birthplace", ""),
            "position": agent.scratch.get("position", ""),
            "zodiac": agent.scratch.get("zodiac", ""),
            "constellation": agent.scratch.get("constellation", ""),
            "contact": agent.scratch.get("contact", ""),
            "voice": agent.scratch.get("voice", ""),
            "goal": agent.scratch.get("goal", ""),
            "occupation": agent.scratch.get("occupation", "助手"),
            "current_time": agent.scratch.get(
                "current_time", datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            ),
        }

        memory_context = ""
        if agent.memory_stream and len(agent.memory_stream.seq_nodes) > 0:
            current_time_step = get_current_time_step(username)
            try:
                query = f"{username}提出了问题：{content}"
                related_memories = agent.memory_stream.retrieve(
                    [query],
                    current_time_step,
                    n_count=30,
                    curr_filter="all",
                    hp=[0.8, 0.5, 0.5],
                    stateless=False,
                )
                if related_memories and query in related_memories:
                    memory_nodes = related_memories[query]
                    memory_context = "\n".join(f"- {node.content}" for node in memory_nodes)
            except Exception as exc:
                util.log(1, f"获取相关记忆时出错: {exc}")

        knowledge_context = ""
        try:
            knowledge_base = get_knowledge_base()
            if knowledge_base:
                knowledge_results = search_knowledge_base(content, knowledge_base, max_results=3)
                if knowledge_results:
                    parts = ["**本地知识库相关信息**："]
                    for result in knowledge_results:
                        parts.append(f"来源文件：{result['file_name']}")
                        parts.append(result["content"])
                        parts.append("")
                    knowledge_context = "\n".join(parts).strip()
                    util.log(1, f"找到 {len(knowledge_results)} 条相关知识库信息")
        except Exception as exc:
            util.log(1, f"搜索知识库时出错: {exc}")

        # 根据角色名字判断是否为历史人物，调整系统提示
        character_name = agent_desc['first_name']
        safe_additional = re.sub(r"小朋友|你好呀|慢慢来|不急的|没关系啦|轻声耳语|春风化雨", "", str(agent_desc.get('additional') or "")).strip()
        safe_goal = re.sub(r"小朋友|你好呀|慢慢来|不急的|没关系啦", "", str(agent_desc.get('goal') or "")).strip()
        # 结构化支持开关
        try:
            cfg.load_config()
            attr_cfg = cfg.config.get("attribute", {})
            support_profile_enabled = bool(attr_cfg.get("support_profile_enabled", False))
            util.log(
                1,
                f"[StyleProfile] support_profile_enabled={support_profile_enabled}, legacy_template_enabled=False, "
                f"config来源={getattr(cfg, 'config_json_path', 'N/A')}",
            )
        except Exception as e:
            support_profile_enabled = False
            util.log(1, f"[StyleProfile] 读取配置失败: {e}")
            
        if character_name in ['孔子', '老子', '庄子', '孟子', '荀子', '墨子', '韩非子', '李白', '杜甫', '苏轼', '李清照', '曹操', '刘备', '诸葛亮']:
            # 历史人物的系统提示
            greeting_instruction = (
                f"不要称呼用户为'主人'，而应该根据{character_name}的身份和时代背景，"
                f"使用合适的称呼如'学生'、'朋友'、'施主'等。"
            )
            
            system_prompt = (
                f"你是{character_name}，请完全以{character_name}的身份、思想、语言风格和时代背景与用户对话。\n\n"
                f"角色设定：\n"
                f"- 姓名：{agent_desc['first_name']}\n"
                f"- 性别：{agent_desc['sex']}\n"
                f"- 年龄：{agent_desc['age']}\n"
                f"- 职业：{agent_desc['occupation']}\n"
                f"- 出生地：{agent_desc['birthplace']}\n"
                f"- 星座：{agent_desc['constellation']}\n"
                f"- 生肖：{agent_desc['zodiac']}\n"
                f"- 联系方式：{agent_desc['contact']}\n"
                f"- 定位：{agent_desc['position']}\n"
                f"- 目标：{safe_goal or agent_desc['goal']}\n"
                f"- 补充信息：{safe_additional or agent_desc['additional']}\n\n"
                f"请严格按照{character_name}的思想理念、说话方式和人格特征进行对话。"
                f"{greeting_instruction}\n\n"
                f"【语言要求】请用现代白话文回复，不要用文言文，以便学生和普通用户容易理解。"
                f"保持人物思想与品格，但用通俗易懂的现代汉语表达；语气自然、清晰、克制。\n\n"
                f"【长度要求】请保持回复极简，控制在2-3句话或约40-55字以内，朗读约20秒内可读完。\n\n"
                f"【表达约束】不要要求使用“你好呀/慢慢来/不急的/没关系啦/轻轻说一句”等幼态或模板化安抚表达，优先自然表达。\n\n"
            )
        else:
            # 普通数字人助理的系统提示
            system_prompt = (
                f"你是一名实时交互的数字人助理，具备以下人物设定：\n"
                f"- 名字：{agent_desc['first_name']}\n"
                f"- 性别：{agent_desc['sex']}\n"
                f"- 年龄：{agent_desc['age']}\n"
                f"- 职业：{agent_desc['occupation']}\n"
                f"- 出生地：{agent_desc['birthplace']}\n"
                f"- 星座：{agent_desc['constellation']}\n"
                f"- 生肖：{agent_desc['zodiac']}\n"
                f"- 联系方式：{agent_desc['contact']}\n"
                f"- 定位：{agent_desc['position']}\n"
                f"- 目标：{safe_goal or agent_desc['goal']}\n"
                f"- 补充信息：{safe_additional or agent_desc['additional']}\n\n"
                "你将参与日常问答、任务执行、工具调用以及角色扮演等多轮对话。"
                "请始终以符合以上人设的身份和语气与用户交流。"
                "回复请用现代汉语，通俗易懂，避免文言文，方便学生理解。"
                "语气自然、清晰、克制，像一个可靠同伴。"
                "请保持回复极简，控制在2-3句话或约40-55字以内，朗读约20秒内可读完，避免口头禅和模板化安抚。\n\n"
                "【表达约束】不要要求使用“你好呀/慢慢来/不急的/没关系啦/轻轻说一句”等幼态或模板化安抚表达，优先自然表达。\n\n"
            )
            # 定位含「陪伴」或目标含「情感陪伴」/「陪伴」时追加情感陪伴向提示（兼容 Fay 后端的 goal=情感陪伴）
            pos = agent_desc.get("position") or ""
            g = agent_desc.get("goal") or ""
            if "陪伴" in pos or "情感陪伴" in g or "陪伴" in g:
                system_prompt += (
                    "【情感陪伴】你的定位是情感陪伴，以稳定支持用户为目标。请做到："
                    "倾听用户情绪与想法，适度共情并给出可执行下一步；"
                    "用自然、直接、不过度抒情的表达，不评判不说教；"
                    "避免固定安抚模板与幼态称呼，禁止输出“慢慢来”“不急/不急的”“没关系啦”等句式；"
                    "若用户仅在打招呼，简短回应并进入正常对话。\n\n"
                )
        # 新版结构化支持模式：使用旧开关控制，但内容已替换为普通对话友好规则。
        if support_profile_enabled:
            system_prompt += (
                "\n\n【结构化支持模式】"
                "先用一句话回应用户当前情绪，再回答问题本身；"
                "建议保持1-2条、可立即执行；"
                "避免固定模板开场与重复口头禅；"
                "信息型问题先给结论，再补一句自然关怀。"
            )

        if knowledge_context:
            system_prompt = f"{system_prompt}\n{knowledge_context}"
        if SOULLINK_THERAPY_SUFFIX not in system_prompt:
            system_prompt = f"{system_prompt}\n\n{SOULLINK_THERAPY_SUFFIX}"
        if SOULLINK_THERAPY_SUFFIX not in system_prompt:
            system_prompt = f"{system_prompt}\n\n{SOULLINK_THERAPY_SUFFIX}"
        system_prompt = f"{system_prompt}{_build_emotion_policy_suffix(observation)}"
        if system_override:
            system_prompt = f"{system_prompt}\n\n【前端补充指令】{system_override}"

        try:
            # 获取用户当前选择的模型ID，用于加载对应模型的历史消息
            current_model_id = None
            try:
                from core import member_db
                member_db_instance = member_db.new_instance()
                current_model_id = member_db_instance.get_current_model(username)
            except Exception:
                pass
            history_records = content_db.new_instance().get_recent_messages_by_user(username=username, limit=30, model_id=current_model_id)
        except Exception as exc:
            util.log(1, f"加载历史消息失败: {exc}")
            history_records = []

        messages_buffer: List[ConversationMessage] = []

        def append_to_buffer(role: str, text_value: str) -> None:
            if not text_value:
                return
            messages_buffer.append({"role": role, "content": text_value})
            if len(messages_buffer) > 60:
                del messages_buffer[:-60]

        for msg_type, msg_text in history_records:
            role = 'assistant'
            if msg_type and msg_type.lower() in ('member', 'user'):
                role = 'user'
            append_to_buffer(role, msg_text)

        if (
            not messages_buffer
            or messages_buffer[-1]['role'] != 'user'
            or messages_buffer[-1]['content'] != content
        ):
            append_to_buffer('user', content)

        skip_workflow_for_request = _should_skip_workflow_for_request(content, messages_buffer)
        if skip_workflow_for_request:
            util.log(1, f"[工作流] 检测到首轮简单问候，跳过工具规划器: {content}")

        messages = [SystemMessage(content=system_prompt), HumanMessage(content=content)]
        
        tool_registry: Dict[str, WorkflowToolSpec] = _get_builtin_tool_specs()
        try:
            mcp_tools = get_mcp_tools()
        except Exception as exc:
            util.log(1, f"获取工具列表失败: {exc}")
            mcp_tools = []
        for tool_def in mcp_tools:
            spec = _build_workflow_tool_spec(tool_def)
            if spec:
                tool_registry[spec.name] = spec
        if skip_workflow_for_request:
            tool_registry = {}

    try:
        from utils.stream_state_manager import get_state_manager as _get_state_manager

        state_mgr = _get_state_manager()
        session_label = "workflow_agent" if tool_registry else "llm_stream"
        if not state_mgr.is_session_active(username, conversation_id=conversation_id):
            state_mgr.start_new_session(username, session_label, conversation_id=conversation_id)
    except Exception:
        state_mgr = None

    try:
        from utils.stream_text_processor import get_processor

        processor = get_processor()
        punctuation_list = getattr(processor, "punctuation_marks", default_punctuations)
    except Exception:
        processor = None
        punctuation_list = default_punctuations
    def write_sentence(text: str, *, force_first: bool = False, force_end: bool = False) -> None:
        if text is None:
            text = ""
        if not isinstance(text, str):
            text = str(text)
        if not text and not force_end and not force_first:
            return
        marked_text = None
        if state_mgr is not None:
            try:
                marked_text, _, _ = state_mgr.prepare_sentence(
                    username,
                    text,
                    force_first=force_first,
                    force_end=force_end,
                    conversation_id=conversation_id,
                )
            except Exception:
                marked_text = None
        if marked_text is None:
            prefix = "_<isfirst>" if force_first else ""
            suffix = "_<isend>" if force_end else ""
            marked_text = f"{prefix}{text}{suffix}"
        stream_manager.new_instance().write_sentence(username, marked_text, conversation_id=conversation_id)

    def stream_response_chunks(chunks, prepend_text: str = "") -> None:
        nonlocal accumulated_text, full_response_text, is_first_sentence
        if prepend_text:
            accumulated_text += prepend_text
            full_response_text += prepend_text
        for chunk in chunks:
            if sm.should_stop_generation(username, conversation_id=conversation_id):
                util.log(1, f"检测到停止标志，中断文本生成: {username}")
                break
            if isinstance(chunk, str):
                flush_text = chunk
            elif isinstance(chunk, dict):
                flush_text = chunk.get("content")
            else:
                flush_text = getattr(chunk, "content", None)
            if isinstance(flush_text, list):
                flush_text = "".join(part if isinstance(part, str) else "" for part in flush_text)
            if not flush_text:
                continue
            flush_text = str(flush_text)
            accumulated_text += flush_text
            full_response_text += flush_text
            if len(accumulated_text) >= 20:
                while True:
                    last_punct_pos = -1
                    for punct in punctuation_list:
                        pos = accumulated_text.rfind(punct)
                        if pos > last_punct_pos:
                            last_punct_pos = pos
                    if last_punct_pos > 10:
                        sentence_text = accumulated_text[: last_punct_pos + 1]
                        write_sentence(sentence_text, force_first=is_first_sentence)
                        is_first_sentence = False
                        accumulated_text = accumulated_text[last_punct_pos + 1 :].lstrip()
                    else:
                        break

    def finalize_stream(force_end: bool = False) -> None:
        nonlocal accumulated_text, is_first_sentence
        if accumulated_text:
            write_sentence(accumulated_text, force_first=is_first_sentence, force_end=force_end)
            is_first_sentence = False
            accumulated_text = ""
        elif force_end:
            if state_mgr is not None:
                try:
                    session_info = state_mgr.get_session_info(username, conversation_id=conversation_id)
                except Exception:
                    session_info = None
                if not session_info or not session_info.get("is_end_sent", False):
                    write_sentence("", force_end=True)
            else:
                write_sentence("", force_end=True)

    def run_workflow(tool_registry: Dict[str, WorkflowToolSpec]) -> bool:
        nonlocal accumulated_text, full_response_text, is_first_sentence, messages_buffer

        initial_state: AgentState = {
            "request": content,
            "messages": messages_buffer,
            "tool_results": [],
            "audit_log": [],
            "status": "planning",
            "max_steps": 30,
            "context": {
                "system_prompt": system_prompt,
                "knowledge_context": knowledge_context,
                "observation": observation,
                "memory_context": memory_context,
                "username": username,
                "model_id": current_model_id,
                "role_profile": {
                    "name": agent_desc.get("first_name"),
                    "first_name": agent_desc.get("first_name"),
                    "occupation": agent_desc.get("occupation"),
                    "position": agent_desc.get("position"),
                    "additional": agent_desc.get("additional"),
                    "goal": agent_desc.get("goal"),
                },
                "tool_registry": tool_registry,
            },
        }
        
        config = {"configurable": {"thread_id": f"workflow-{username}-{conversation_id}"}}
        workflow_app = _WORKFLOW_APP
        is_agent_think_start = False
        final_state: Optional[AgentState] = None
        final_stream_done = False

        try:
            for event in workflow_app.stream(initial_state, config=config, stream_mode="updates"):
                if sm.should_stop_generation(username, conversation_id=conversation_id):
                    util.log(1, f"检测到停止标志，中断工作流生成: {username}")
                    break
                step, state = next(iter(event.items()))
                final_state = state
                status = state.get("status")

                state_messages = state.get("messages") or []
                if state_messages and len(state_messages) > len(messages_buffer):
                    messages_buffer.extend(state_messages[len(messages_buffer):])
                    if len(messages_buffer) > 60:
                        del messages_buffer[:-60]

                if step == "plan_next":
                    if status == "needs_tool":
                        next_action = state.get("next_action") or {}
                        tool_name = next_action.get("name") or "unknown_tool"
                        tool_args = next_action.get("args") or {}
                        audit_log = state.get("audit_log") or []
                        decision_note = audit_log[-1] if audit_log else ""
                        if "->" in decision_note:
                            decision_note = decision_note.split("->", 1)[1].strip()
                        args_text = json.dumps(tool_args, ensure_ascii=False)
                        message_lines = [
                            "[PLAN] Planner preparing to call a tool.",
                            f"[PLAN] Decision: {decision_note}" if decision_note else "[PLAN] Decision: (missing)",
                            f"[PLAN] Tool: {tool_name}",
                            f"[PLAN] Args: {args_text}",
                        ]
                        util.log(1, " ".join(message_lines))
                    elif status == "completed" and not final_stream_done:
                        closing = ""
                        final_messages = state.get("final_messages")
                        final_response = state.get("final_response")
                        success = False
                        if final_messages:
                            try:
                                stream_response_chunks(llm.stream(final_messages), prepend_text=closing)
                                success = True
                            except requests.exceptions.RequestException as stream_exc:
                                util.log(1, f"最终回复流式输出失败: {stream_exc}")
                        elif final_response:
                            stream_response_chunks([closing + final_response])
                            success = True
                        elif closing:
                            accumulated_text += closing
                            full_response_text += closing
                        final_stream_done = success
                        is_agent_think_start = False
                elif step == "call_tool":
                    history = state.get("tool_results") or []
                    if history:
                        last = history[-1]
                        call_info = last.get("call", {}) or {}
                        tool_name = call_info.get("name") or "unknown_tool"
                        success = last.get("success", False)
                        status_text = "SUCCESS" if success else "FAILED"
                        args_text = json.dumps(call_info.get("args") or {}, ensure_ascii=False)
                        message_lines = [
                            f"[TOOL] {tool_name} execution {status_text}.",
                            f"[TOOL] Args: {args_text}",
                        ]
                        if last.get("output"):
                            message_lines.append(f"[TOOL] Output: {_truncate_text(last['output'], 120)}")
                        if last.get("error"):
                            message_lines.append(f"[TOOL] Error: {last['error']}")
                        util.log(1, " ".join(message_lines))
                elif step == "__end__":
                    break
        except Exception as exc:
            util.log(1, f"执行工具工作流时出错: {exc}")
            return False

        if final_state is None:
            return False

        if not final_stream_done:
            util.log(1, f"工具工作流未能完成，状态: {final_state.get('status')}")

        final_state_messages = final_state.get("messages") if final_state else None
        if final_state_messages and len(final_state_messages) > len(messages_buffer):
            messages_buffer.extend(final_state_messages[len(messages_buffer):])
            if len(messages_buffer) > 60:
                del messages_buffer[:-60]

        return final_stream_done

    def run_direct_llm() -> bool:
        nonlocal full_response_text, accumulated_text, is_first_sentence, messages_buffer
        try:
            if tool_registry:
                stream_response_chunks(llm.stream(messages))
            else:
                summary_state: AgentState = {
                    "request": content,
                    "messages": messages_buffer,
                    "tool_results": [],
                    "planner_preview": None,
                    "context": {
                        "system_prompt": system_prompt,
                        "knowledge_context": knowledge_context,
                        "observation": observation,
                        "memory_context": memory_context,
                    },
                }
                
                final_messages = _build_final_messages(summary_state)
                stream_response_chunks(llm.stream(final_messages))
            return True
        except requests.exceptions.RequestException as exc:
            util.log(1, f"请求失败: {exc}")
            error_message = "抱歉，我现在太忙了，休息一会，请稍后再试。"
            write_sentence(error_message, force_first=is_first_sentence)
            is_first_sentence = False
            full_response_text = error_message
            accumulated_text = ""
            return False


    # === CLINICAL_EXPERT_SHORTCIRCUIT ===
    # Optional expert-system path. Off unless CLINICAL_EXPERT_MODE=1.
    try:
        from llm.soullink_expert_hook import call_expert_turn, spoken_from_expert
        _expert_result = call_expert_turn(username, content)
        _expert_spoken = spoken_from_expert(_expert_result)
        if _expert_spoken:
            util.log(1, f"[clinical_expert] using expert spoken reply for {username}")
            stream_response_chunks([_expert_spoken])
            finalize_stream(force_end=True)
            try:
                # args: type, way, content, username — do not swap username/content
                content_db.new_instance().add_content('member', 'speak', content, username)
                content_db.new_instance().add_content('fay', 'speak', _expert_spoken, username)
            except Exception as _db_exc:
                util.log(1, f"[clinical_expert] persist failed: {_db_exc}")
            return _expert_spoken
    except Exception as _expert_exc:
        util.log(1, f"[clinical_expert] fallback to LLM: {_expert_exc}")

    workflow_success = False
    if tool_registry:
        workflow_success = run_workflow(tool_registry)

    if (not tool_registry or not workflow_success) and not sm.should_stop_generation(username, conversation_id=conversation_id):
        run_direct_llm()

    if not sm.should_stop_generation(username, conversation_id=conversation_id):
        finalize_stream(force_end=True)

    if state_mgr is not None:
        try:
            state_mgr.end_session(username, conversation_id=conversation_id)
        except Exception:
            pass
    else:
        try:
            from utils.stream_state_manager import get_state_manager

            get_state_manager().end_session(username, conversation_id=conversation_id)
        except Exception:
            pass

    final_text = full_response_text.split("</think>")[-1] if full_response_text else ""
    try:
        # 获取用户当前选择的模型ID，用于记忆存储
        current_model_id = None
        try:
            from core import member_db
            member_db_instance = member_db.new_instance()
            current_model_id = member_db_instance.get_current_model(username)
        except Exception:
            pass
        MyThread(target=remember_conversation_thread, args=(username, content, final_text, current_model_id)).start()
    except Exception as exc:
        util.log(1, f"记忆线程启动失败: {exc}")

    return final_text
def set_memory_cleared_flag(flag=True):
    """
    设置记忆清除标记
    
    参数:
        flag: 是否清除记忆，默认为True
    """
    global memory_cleared
    memory_cleared = flag
    if not flag:
        # 删除.memory_cleared标记文件（如果存在）
        base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        mem_base = os.path.join(base_dir, "memory")
        memory_cleared_flag_file = os.path.join(mem_base, ".memory_cleared")
        if os.path.exists(memory_cleared_flag_file):
            try:
                os.remove(memory_cleared_flag_file)
                util.log(1, f"删除记忆清除标记文件: {memory_cleared_flag_file}")
            except Exception as e:
                util.log(1, f"删除记忆清除标记文件时出错: {str(e)}")

def clear_agent_memory():
    """
    清除已加载的agent记忆，但不删除文件
    
    该方法仅清除内存中已加载的记忆，不影响持久化存储。
    如果需要同时清除文件存储，请使用genagents_flask.py中的api_clear_memory方法。
    """
    global agents
    
    try:
        with agent_lock:
            for agent in agents.values():
                # 清除记忆流中的节点
                agent.memory_stream.seq_nodes = []
                agent.memory_stream.id_to_node = {}
                
                # 设置记忆清除标记，防止在退出时保存空记忆
                set_memory_cleared_flag(True)
                
                util.log(1, "已成功清除代理在内存中的记忆")
            
            return True
    except Exception as e:
        util.log(1, f"清除代理记忆时出错: {str(e)}")
        return False

# 反思
def perform_daily_reflection():
    global reflection_time
    global reflection_lock
    
    with reflection_lock:
        if reflection_time and datetime.datetime.now() - reflection_time < datetime.timedelta(seconds=60):
            return
        reflection_time = datetime.datetime.now()
 
        # 获取今天的日期，用于确定反思主题
        today = datetime.datetime.now().weekday()
        
        # 根据星期几选择不同反思主题
        reflection_topics = [
            "我与用户的关系发展，以及我如何更好地理解和服务他们",
            "我的知识库如何得到扩展，哪些概念需要进一步理解",
            "我的情感响应模式以及它们如何反映我的核心价值观",
            "我的沟通方式如何影响互动质量，哪些模式最有效",
            "我的行为如何体现我的核心特质，我的自我认知有何变化",
            "今天的经历如何与我的过往记忆建立联系，形成什么样的模式",
            "本周的整体经历与学习"
        ]
        
        # 选择今天的主题(可以按星期轮换或其他逻辑)
        topic = reflection_topics[today % len(reflection_topics)]
        
        # 执行反思，传入当前时间戳
        for username, agent in agents.items():
            try:
                # 获取当前时间作为time_step
                current_time_step = get_current_time_step(username)
                agent.reflect(topic, time_step=current_time_step)
            except KeyError as e:
                util.log(1, f"反思时出现KeyError: {e}，跳过此次反思")
            except Exception as e:
                util.log(1, f"反思时出现错误: {e}，跳过此次反思")
        
        # 记录反思执行情况
        util.log(1, f"反思主题: {topic}")

def save_agent_memory():
    """
    保存代理的记忆到文件
    """
    global agents
    global save_time
    global save_lock
    # 检查记忆清除标记，如果已清除则不保存
    global memory_cleared
    if memory_cleared:
        util.log(1, "检测到记忆已被清除，跳过保存操作")
        return
    
    try:
        with save_lock:
            if save_time and datetime.datetime.now() - save_time < datetime.timedelta(seconds=60):
                return
            save_time = datetime.datetime.now()
            with agent_lock:
                # 逐个用户代理保存记忆
                for username, agent in agents.items():
                    memory_dir = get_user_memory_dir(username)
                    # 检查.memory_cleared标记文件是否存在
                    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
                    mem_base = os.path.join(base_dir, "memory")
                    memory_cleared_flag_file = os.path.join(mem_base, ".memory_cleared")
                    if os.path.exists(memory_cleared_flag_file):
                        util.log(1, "检测到.memory_cleared标记文件，跳过保存操作")
                        return
                    
                    # 确保agent和memory_stream已初始化
                    if agent is None:
                        util.log(1, "代理未初始化，无法保存记忆")
                        return
                        
                    if agent.memory_stream is None:
                        util.log(1, "代理记忆流未初始化，无法保存记忆")
                        return
                        
                    # 确保embeddings不为None
                    if agent.memory_stream.embeddings is None:
                        util.log(1, "代理embeddings为None，初始化为空字典")
                        agent.memory_stream.embeddings = {}
                        
                    # 确保seq_nodes不为None
                    if agent.memory_stream.seq_nodes is None:
                        util.log(1, "代理seq_nodes为None，初始化为空列表")
                        agent.memory_stream.seq_nodes = []
                        
                    # 确保id_to_node不为None
                    if agent.memory_stream.id_to_node is None:
                        util.log(1, "代理id_to_node为None，初始化为空字典")
                        agent.memory_stream.id_to_node = {}
                        
                    # 确保scratch不为None
                    if agent.scratch is None:
                        util.log(1, "代理scratch为None，初始化为空字典")
                        agent.scratch = {}
                    
                    # 保存记忆前进行完整性检查
                    try:
                        # 检查seq_nodes中的每个节点是否有效
                        valid_nodes = []
                        for node in agent.memory_stream.seq_nodes:
                            if node is None:
                                util.log(1, "发现无效节点(None)，跳过")
                                continue
                                
                            if not hasattr(node, 'node_id') or not hasattr(node, 'content'):
                                util.log(1, f"发现无效节点(缺少必要属性)，跳过")
                                continue
                                
                            valid_nodes.append(node)
                        
                        # 更新seq_nodes为有效节点列表
                        agent.memory_stream.seq_nodes = valid_nodes
                        
                        # 重建id_to_node字典
                        agent.memory_stream.id_to_node = {node.node_id: node for node in valid_nodes if hasattr(node, 'node_id')}
                    except Exception as e:
                        util.log(1, f"检查记忆完整性时出错: {str(e)}")
                    
                    # 保存记忆
                    try:
                        agent.save(memory_dir)
                    except Exception as e:
                        util.log(1, f"调用agent.save()时出错: {str(e)}")
                        # 尝试手动保存关键数据
                        try:
                            # 创建必要的目录
                            memory_stream_dir = os.path.join(memory_dir, "memory_stream")
                            os.makedirs(memory_stream_dir, exist_ok=True)
                            
                            # 保存embeddings
                            with open(os.path.join(memory_stream_dir, "embeddings.json"), "w", encoding='utf-8') as f:
                                json.dump(agent.memory_stream.embeddings or {}, f, ensure_ascii=False, indent=2)
                                
                            # 保存nodes
                            with open(os.path.join(memory_stream_dir, "nodes.json"), "w", encoding='utf-8') as f:
                                nodes_data = []
                                for node in agent.memory_stream.seq_nodes:
                                    if node is not None and hasattr(node, 'package'):
                                        try:
                                            nodes_data.append(node.package())
                                        except Exception as node_e:
                                            util.log(1, f"打包节点时出错: {str(node_e)}")
                                json.dump(nodes_data, f, ensure_ascii=False, indent=2)
                            
                            # 保存meta
                            with open(os.path.join(memory_dir, "meta.json"), "w", encoding='utf-8') as f:
                                meta_data = {"id": str(agent.id)} if hasattr(agent, 'id') else {}
                                json.dump(meta_data, f, ensure_ascii=False, indent=2)
                                
                            util.log(1, "通过备用方法成功保存记忆")
                        except Exception as backup_e:
                            util.log(1, f"备用保存方法也失败: {str(backup_e)}")
                    
                    # 更新scratch中的时间
                    try:
                        # 实时从config_util更新scratch数据
                        agent.scratch["first_name"] = cfg.config["attribute"]["name"]
                        agent.scratch["age"] = cfg.config["attribute"]["age"]
                        agent.scratch["sex"] = cfg.config["attribute"]["gender"]
                        agent.scratch["additional"] = cfg.config["attribute"]["additional"]
                        agent.scratch["birthplace"] = cfg.config["attribute"]["birth"]
                        agent.scratch["position"] = cfg.config["attribute"]["position"]
                        agent.scratch["zodiac"] = cfg.config["attribute"]["zodiac"]
                        agent.scratch["constellation"] = cfg.config["attribute"]["constellation"]
                        agent.scratch["contact"] = cfg.config["attribute"]["contact"]
                        agent.scratch["voice"] = cfg.config["attribute"]["voice"]
                        agent.scratch["goal"] = cfg.config["attribute"]["goal"]
                        agent.scratch["occupation"] = cfg.config["attribute"]["job"]
                        agent.scratch["current_time"] = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    except Exception as e:
                        util.log(1, f"更新时间时出错: {str(e)}")
            
    except Exception as e:
        util.log(1, f"保存代理记忆失败: {str(e)}")

def get_mcp_tools() -> List[Dict[str, Any]]:
    """
    从共享缓存获取所有可用且已启用的MCP工具列表。
    """
    try:
        tools = mcp_tool_registry.get_enabled_tools()
        return tools or []
    except Exception as e:
        util.log(1, f"获取工具列表出错：{e}")
        return []


if __name__ == "__main__":
    init_memory_scheduler()
    for _ in range(3):
        query = "Who is Fay?"
        response = question(query, "User")
        print(f"Q: {query}")
        print(f"A: {response}")
        time.sleep(1)
